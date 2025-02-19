"""
Document Processing Service for the RAG Pipeline.

This module handles the extraction and processing of text from various document formats:
- PDF files (using PDFReader from llama_index)
- Excel spreadsheets (using pandas)
- PowerPoint presentations (using python-pptx)
- Word documents (using python-docx)

Key Features:
1. Text extraction from multiple document formats
2. Document chunking with configurable size and overlap
3. Metadata preservation and enhancement
4. Unique document ID assignment
5. Error handling and logging

The processing follows these general steps:
1. Load document from file
2. Extract text content
3. Split into manageable chunks
4. Assign metadata (source, IDs, etc.)
5. Create Document objects for indexing

Note: This service is part of the data preparation pipeline and feeds
into the indexing service for search index creation.
"""

import logging  # Python's built-in logging facility
from pathlib import Path  # Object-oriented filesystem paths
from typing import List, Dict  # Type hints for better code clarity
import os  # Operating system interface for file operations

# Document processing libraries
from llama_index.core import Document  # Base document class for RAG
from llama_index.readers.file import PDFReader  # PDF text extraction
import pandas as pd  # Data manipulation library for Excel files
from docx import Document as WordDocument  # Word document processing
from pptx import Presentation  # PowerPoint processing
import uuid  # Unique identifier generation

# Configure module logger
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """
    Service for processing various document types into indexable chunks.
    
    This class handles:
    1. Document loading from multiple formats
    2. Text extraction and cleaning
    3. Content chunking with overlap
    4. Metadata assignment
    
    The processor supports:
    - PDF files: Full document with page-based chunking
    - Excel files: Row-based processing with column preservation
    - PowerPoint: Slide-based extraction
    - Word documents: Heading-based chunking with overlap
    
    Each document is processed into Document objects that contain:
    - Extracted and cleaned text
    - Source metadata (filename, page/slide numbers, etc.)
    - Unique document ID
    - Chunk information (for split documents)
    """

    def __init__(self, max_tokens: int = 1024, overlap_tokens: int = 50):
        """
        Initialize document processor with chunking parameters.
        
        The processor splits large documents into smaller chunks for better
        search and retrieval. Overlapping tokens help maintain context
        across chunk boundaries.
        
        Args:
            max_tokens (int): Maximum number of tokens per chunk
                Default is 1024, which is suitable for most LLMs
            overlap_tokens (int): Number of overlapping tokens between chunks
                Default is 50, which helps maintain context between chunks
                
        Note: Current implementation uses character-based chunking as an
        approximation. A proper tokenizer should be used in production.
        """
        self.max_tokens = max_tokens
        self.overlap_tokens = overlap_tokens
        self.pdf_reader = PDFReader()

    def load_pdf_docs(self, pdf_paths: List[Path]) -> List[Document]:
        """
        Load and process PDF documents into searchable chunks.
        
        This method:
        1. Validates each PDF file (existence and readability)
        2. Extracts text content using PDFReader
        3. Assigns unique IDs and metadata
        4. Handles memory cleanup after processing
        
        The PDFReader automatically handles:
        - Page-based chunking
        - Text extraction from various PDF formats
        - Basic structure preservation
        
        Args:
            pdf_paths (List[Path]): List of paths to PDF files to process
            
        Returns:
            List[Document]: List of processed Document objects, one per page
            
        Error Handling:
        - Skips files that don't exist or aren't readable
        - Logs warnings for PDFs with no extractable content
        - Continues processing remaining files if one fails
        - Ensures cleanup of file handles
        """
        all_docs = []
        for pdf_path in pdf_paths:
            try:
                # Log the start of processing for this PDF
                logger.info(f"Starting to process PDF: {pdf_path.name} (size: {pdf_path.stat().st_size} bytes)")
                
                # Check if file exists and is readable
                if not pdf_path.exists():
                    logger.error(f"PDF file not found: {pdf_path}")
                    continue
                    
                if not os.access(pdf_path, os.R_OK):
                    logger.error(f"PDF file not readable: {pdf_path}")
                    continue
                
                # Load and process the PDF
                docs = []
                try:
                    docs = self.pdf_reader.load_data(str(pdf_path))
                finally:
                    # Ensure any file handles are closed
                    import gc
                    gc.collect()
                
                if not docs:
                    logger.warning(f"No content extracted from PDF: {pdf_path.name}")
                    continue
                    
                # Assign document IDs and metadata
                docs = self._assign_doc_ids(docs)
                for doc in docs:
                    doc = self._extract_metadata(doc, source_file=pdf_path.name)
                    
                # Add to the collection
                all_docs.extend(docs)
                logger.info(f"Successfully processed PDF: {pdf_path.name} - extracted {len(docs)} pages")
                
            except Exception as e:
                logger.error(f"Failed to process PDF {pdf_path.name}: {str(e)}")
                logger.exception("Detailed error information:")
                
        logger.info(f"Completed processing {len(pdf_paths)} PDF files, extracted {len(all_docs)} total pages")
        return all_docs

    def load_excel_docs(self, excel_paths: List[Path]) -> List[Document]:
        """
        Load and process Excel spreadsheets into searchable documents.
        
        This method processes Excel files by:
        1. Reading each sheet in the workbook
        2. Converting each row into a document
        3. Preserving column names as field labels
        4. Adding sheet and row metadata
        
        The processing approach:
        - Each row becomes a separate document
        - Column headers are preserved in the text: "column_name: value"
        - Sheet name and row number are stored in metadata
        
        Args:
            excel_paths (List[Path]): List of paths to Excel files
            
        Returns:
            List[Document]: List of processed Document objects, one per row
            
        Example:
            For a spreadsheet with columns "Name" and "Age":
            Text: "name: John Smith, age: 30"
            Metadata: {"sheet": "Sheet1", "row_number": 1, "source": "data.xlsx"}
        """
        all_docs = []
        for excel_path in excel_paths:
            try:
                # Open Excel file and process each sheet
                xls = pd.ExcelFile(excel_path)
                for sheet_name in xls.sheet_names:
                    # Read sheet into DataFrame
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    
                    # Process each row
                    for idx, row in df.iterrows():
                        # Convert row to text with column labels
                        row_dict = ", ".join(f"{col.strip().lower()}: {str(row[col]).strip()}" 
                                           for col in df.columns)
                        
                        # Create metadata
                        metadata = {
                            "sheet": sheet_name,
                            "row_number": idx + 1,
                        }
                        
                        # Create and store document
                        doc = Document(text=row_dict, metadata=metadata)
                        all_docs.append(doc)
                        
                logger.info(f"Processed Excel: {excel_path.name} with {len(df)} rows.")
                
            except Exception as e:
                logger.error(f"Failed to process Excel {excel_path.name}: {e}")
        
        # Assign IDs and metadata to all documents
        all_docs = self._assign_doc_ids(all_docs)
        for doc in all_docs:
            doc = self._extract_metadata(doc, source_file=excel_path.name)
            
        return all_docs

    def load_ppt_docs(self, ppt_paths: List[Path]) -> List[Document]:
        """
        Load and process PowerPoint presentations into searchable documents.
        
        This method processes presentations by:
        1. Extracting text from each slide
        2. Preserving text from all shapes (text boxes, tables, etc.)
        3. Maintaining slide numbers in metadata
        
        Processing approach:
        - Each slide becomes a separate document
        - Text from all shapes is combined with newlines
        - Slide numbers are preserved in metadata
        
        Args:
            ppt_paths (List[Path]): List of paths to PowerPoint files
            
        Returns:
            List[Document]: List of processed Document objects, one per slide
            
        Note:
            - Handles both .ppt and .pptx formats
            - Extracts text from various shape types (if they have text)
            - Preserves slide order through metadata
        """
        all_docs = []
        for ppt_path in ppt_paths:
            try:
                # Load presentation and process each slide
                prs = Presentation(ppt_path)
                for idx, slide in enumerate(prs.slides):
                    # Extract text from all shapes that have text
                    slide_text = "\n".join(shape.text for shape in slide.shapes 
                                         if hasattr(shape, "text"))
                    
                    # Create metadata with slide number
                    metadata = {
                        "slide_number": idx + 1,
                    }
                    
                    # Create and store document
                    doc = Document(text=slide_text, metadata=metadata)
                    all_docs.append(doc)
                    
                logger.info(f"Processed PowerPoint: {ppt_path.name} with {len(prs.slides)} slides.")
                
            except Exception as e:
                logger.error(f"Failed to process PowerPoint {ppt_path.name}: {e}")
        
        # Assign IDs and metadata to all documents
        all_docs = self._assign_doc_ids(all_docs)
        for doc in all_docs:
            doc = self._extract_metadata(doc, source_file=ppt_path.name)
            
        return all_docs

    def load_doc_docs(self, doc_paths: List[Path]) -> List[Document]:
        """
        Load and process Word documents into searchable chunks.
        
        This method processes Word documents by:
        1. Using headings to identify logical sections
        2. Chunking large sections for better processing
        3. Preserving document structure in metadata
        
        Processing approach:
        - Documents are split at heading boundaries
        - Large sections are further split with overlap
        - Heading text is preserved in chunk metadata
        
        Args:
            doc_paths (List[Path]): List of paths to Word documents
            
        Returns:
            List[Document]: List of processed Document objects
            
        Note:
            - Handles both .doc and .docx formats
            - Preserves document structure through headings
            - Uses overlapping chunks for large sections
            - Maintains heading hierarchy in metadata
        """
        all_docs = []
        for doc_path in doc_paths:
            try:
                # Load and process Word document
                word_doc = WordDocument(doc_path)
                
                # Split document into chunks based on headings
                docs = self._chunk_by_headings(word_doc)
                
                # Add source metadata to each chunk
                for doc in docs:
                    doc = self._extract_metadata(doc, source_file=doc_path.name)
                    
                # Add to collection
                all_docs.extend(docs)
                logger.info(f"Processed Word Document: {doc_path.name} with {len(docs)} chunks.")
                
            except Exception as e:
                logger.error(f"Failed to process Word Document {doc_path.name}: {e}")
                
        return all_docs

    def _chunk_by_headings(self, word_doc: WordDocument) -> List[Document]:
        """
        Split a Word document into chunks based on heading structure.
        
        This method:
        1. Identifies heading paragraphs in the document
        2. Uses headings as natural section boundaries
        3. Groups content between headings
        4. Splits large sections if needed
        
        The approach preserves document structure by:
        - Keeping content under its relevant heading
        - Maintaining heading hierarchy
        - Preserving section context in metadata
        
        Args:
            word_doc (WordDocument): Loaded Word document to process
            
        Returns:
            List[Document]: List of document chunks with metadata
            
        Note:
            - Recognizes all heading levels (Heading 1-9)
            - Large sections are further split using _split_large_section
            - Each chunk includes its heading in metadata
        """
        current_section = []  # Accumulate paragraphs in current section
        current_metadata = None  # Metadata for current section
        all_docs = []  # All processed document chunks

        # Process each paragraph
        for paragraph in word_doc.paragraphs:
            # Check if this is a heading paragraph
            if paragraph.style.name.startswith("Heading"):
                # Process previous section if it exists
                if current_section:
                    docs = self._split_large_section(current_section, current_metadata)
                    all_docs.extend(docs)
                    current_section = []
                    
                # Start new section with this heading
                current_metadata = {"section": paragraph.text.strip()}
                
            # Add paragraph to current section
            current_section.append(paragraph.text.strip())

        # Process final section if it exists
        if current_section:
            docs = self._split_large_section(current_section, current_metadata)
            all_docs.extend(docs)

        # Assign unique IDs to all chunks
        return self._assign_doc_ids(all_docs)

    def _split_large_section(self, section: List[str], metadata: Dict) -> List[Document]:
        """
        Split a large text section into smaller, overlapping chunks.
        
        This method implements a simple chunking strategy:
        1. Joins text with newlines
        2. Splits into chunks of approximately max_tokens size
        3. Ensures chunks break at word boundaries
        4. Creates overlap between chunks
        
        The chunking approach:
        - Uses character count as a proxy for tokens
        - Breaks chunks at word boundaries to avoid mid-word splits
        - Creates overlap between chunks for context preservation
        - Preserves original metadata in all chunks
        
        Args:
            section (List[str]): List of text paragraphs to chunk
            metadata (Dict): Metadata to preserve in all chunks
            
        Returns:
            List[Document]: List of overlapping document chunks
            
        Note:
            This is a simplified implementation using characters instead of
            actual tokens. In production, you should use a proper tokenizer
            for more accurate chunk sizes.
        """
        # Join paragraphs with newlines
        text = "\n".join(section)
        chunk_size = self.max_tokens
        overlap = self.overlap_tokens
        
        # Create chunks with overlap
        chunks = []
        start = 0
        while start < len(text):
            # Calculate end position for this chunk
            end = start + chunk_size
            if end < len(text):
                # Find the last word boundary before end
                end = text.rfind(" ", start, end)
                
            # Extract chunk
            chunk = text[start:end]
            
            # Create metadata for this chunk
            chunk_metadata = {
                **(metadata or {}),  # Include original metadata
                "chunk_number": len(chunks) + 1,  # Add chunk number
                "doc_id": str(uuid.uuid4()),  # Unique ID for chunk
            }
            
            # Create and store chunk document
            chunks.append(Document(text=chunk, metadata=chunk_metadata))
            
            # Move start position for next chunk, including overlap
            start = end - overlap if end < len(text) else len(text)
        
        return chunks

    @staticmethod
    def _assign_doc_ids(documents: List[Document]) -> List[Document]:
        """
        Assign unique identifiers to a list of documents.
        
        This method:
        1. Generates a UUID for each document
        2. Stores it in the document's metadata
        3. Returns the modified document list
        
        The UUID ensures:
        - Globally unique identification
        - Consistent reference tracking
        - Safe document deduplication
        
        Args:
            documents (List[Document]): List of documents needing IDs
            
        Returns:
            List[Document]: Same documents with added IDs in metadata
            
        Note:
            This method modifies the documents in place but also returns
            them for convenience in method chaining.
        """
        for doc in documents:
            doc.metadata["doc_id"] = str(uuid.uuid4())
        return documents

    @staticmethod
    def _extract_metadata(doc: Document, source_file: str) -> Document:
        """
        Extract and assign source metadata to a document.
        
        This method:
        1. Adds source filename to document metadata
        2. Preserves existing metadata
        3. Returns the enhanced document
        
        The source tracking enables:
        - Origin tracing for each chunk
        - Reference back to original documents
        - User-friendly source citations
        
        Args:
            doc (Document): Document to enhance with metadata
            source_file (str): Name of the source file
            
        Returns:
            Document: Same document with added source metadata
            
        Note:
            This method modifies the document in place but also returns
            it for convenience in method chaining.
        """
        doc.metadata["source"] = source_file
        return doc 