# data_preparation/data_preparation.py


########################################
# IMPORT DEPENDENCIES
########################################

import os
import json
import logging
import uuid
import textwrap
from pathlib import Path
from typing import List, Set
# filepath: /c:/Users/admin/Desktop/RAG-pipeline-develop/RAG-pipeline-develop/data_preparation/data_preparation.py
from dotenv import load_dotenv

# Load environment variables from .env (for local development)
load_dotenv(dotenv_path=Path(".env"))

# Access environment variables
api_key = os.getenv("API_KEY")
api_base = os.getenv("API_BASE")
if not api_key or not api_base:
    raise EnvironmentError("API_KEY and API_BASE must be set as environment variables.")

import httpx
import faiss
import openai
import pandas as pd
import pdfplumber
from docx import Document as WordDocument
from pptx import Presentation
from transformers import GPT2TokenizerFast
from whoosh.index import create_in, open_dir, exists_in
from whoosh.fields import Schema, TEXT, ID, NUMERIC

from fastapi import FastAPI, UploadFile, HTTPException
from fastapi.responses import JSONResponse
import uvicorn
import gradio as gr
from logging.handlers import RotatingFileHandler

from llama_index.core import Document
from llama_index.readers.file import PDFReader
from llama_index.embeddings.nebius import NebiusEmbedding
from llama_index.vector_stores.faiss import FaissVectorStore
from llama_index.core import Settings, VectorStoreIndex, StorageContext, load_index_from_storage
from llama_index.core.storage.docstore.simple_docstore import SimpleDocumentStore
from llama_index.core.storage.index_store import SimpleIndexStore

# from dotenv import load_dotenv

import threading

########################################
# INITIALIZE FASTAPI APP
########################################

app = FastAPI(title="Data Preparation Pipeline API")

########################################
# LOAD ENVIRONMENT VARIABLES
########################################

# Load environment variables from .env (for local development)
# load_dotenv(dotenv_path=Path("config/.env"))

########################################
# CONFIG & INITIAL SETUP
########################################

# Load configuration
CONFIG_PATH = Path("config/config.sample.json")
with open(CONFIG_PATH, encoding="utf-8-sig") as config_file:
    config = json.load(config_file)

# Access environment variables
api_key = os.getenv("API_KEY")
api_base = os.getenv("API_BASE")
if not api_key or not api_base:
    raise EnvironmentError("API_KEY and API_BASE must be set as environment variables.")

# Set OpenAI API credentials
openai.api_key = api_key
openai.api_base = api_base

########################################
# DEFINE CHUNK PARAMETERS
########################################

max_tokens = config['chunking']['max_tokens']
overlap_tokens = config['chunking']['overlap_tokens']

########################################
# DEFINE DIRECTORIES
########################################

data_directory = Path(config["directories"]["data_directory"])  # 'data'
to_process_dir = Path(config["directories"]["to_process_dir"])  # 'data/to_process'
processed_dir = Path(config["directories"]["processed_dir"])    # 'data/processed'

# Index directories
whoosh_index_path = Path(config["directories"]["whoosh_index_path"])  # 'data/whoosh_index'
faiss_index_path = Path(config["directories"]["faiss_index_path"])    # 'data/faiss_index'

# Log directories
inference_log_folder = Path(config["directories"]["inference_log_folder"])  # 'logs/inference'
data_preparation_log_folder = Path(config["directories"]["data_preparation_log_folder"])  # 'logs/data_preparation'

# Ensure necessary directories exist
for folder in [
    to_process_dir, processed_dir, whoosh_index_path, faiss_index_path,
    inference_log_folder, data_preparation_log_folder
]:
    folder.mkdir(parents=True, exist_ok=True)

########################################
# SETUP LOGGING
########################################

def setup_logging(
    log_folder: Path,
    log_file: str = "data_preparation.log",
    max_bytes: int = 5 * 1024 * 1024,  # 5 MB
    backup_count: int = 3,
) -> None:
    os.makedirs(log_folder, exist_ok=True)
    log_file_path = log_folder / log_file

    rotating_handler = RotatingFileHandler(
        filename=log_file_path,
        maxBytes=max_bytes,
        backupCount=backup_count,
        encoding="utf-8",
    )

    logging.basicConfig(
        level=config['logging']['level'],
        format=config['logging']['format'],
        handlers=[
            rotating_handler,
            logging.StreamHandler(),
        ],
        force=True,
    )

    logging.info("Logging with rotation has been successfully configured.")

# Initialize logging
setup_logging(
    log_folder=data_preparation_log_folder,
    log_file=config['logging']['log_file'],
    max_bytes=config['logging']['max_bytes'],
    backup_count=config['logging']['backup_count'],
)

########################################
# LOAD EMBEDDING MODEL
########################################

# Initialize tokenizer
tokenizer = GPT2TokenizerFast.from_pretrained("gpt2")

# Setup NebiusEmbedding model
custom_http_client = httpx.Client(timeout=60.0)
embedding_model = NebiusEmbedding(
    api_key=api_key,
    model_name=config['embedding_model']['model_name'],  # e.g., "BAAI/bge-en-icl"
    http_client=custom_http_client,
    api_base=api_base,
)

# Determine embedding dimension
sample_text = "test text to create an embedding"
sample_embedding = embedding_model.get_text_embedding(sample_text)
embedding_dimension = len(sample_embedding)

logging.info(f"\nEmbedding Model {config['embedding_model']['model_name']} creates {embedding_dimension}-dimensional vectors.")
logging.info(f"Sample embedding for text '{sample_text}': {sample_embedding[:10]} ...\n")

# Set the embedding model in LlamaIndex settings
Settings.embed_model = embedding_model

########################################
# DEFINE UTILITIES FOR PROCESSING
########################################

def assign_doc_id(documents: list) -> list:
    """
    Assign a unique UUID to each document.
    """
    for doc in documents:
        doc_id = str(uuid.uuid4())
        doc.metadata["doc_id"] = doc_id
    return documents

def extract_metadata(doc: Document, source_file: str, additional_info: dict) -> Document:
    """
    Extract and assign metadata to a document.
    """
    doc.metadata["source"] = source_file
    for key, value in additional_info.items():
        doc.metadata[key] = value
    return doc

def log_new_files(file_list: List[Path], file_type: str) -> None:
    """
    Log and print new files found for processing.
    """
    if not file_list:
        logging.info(f"No new {file_type.upper()} files found in {to_process_dir}.")
        print(f"No new {file_type.upper()} files to process.\n")
    else:
        logging.info(f"Found {len(file_list)} new {file_type.upper()} files to process.")
        print(f"New {file_type.upper()} Files Found:")
        for file in file_list:
            print(file)
        print()

def load_processed_files(record_path: Path) -> Set[str]:
    """
    Load the set of already processed file names.
    """
    if record_path.exists():
        with open(record_path, 'r') as f:
            processed = set(json.load(f))
        logging.info(f"Loaded processed files from {record_path}.")
    else:
        processed = set()
        logging.info(f"No processed files record found at {record_path}. Starting fresh.")
    return processed

def move_files_to_processed(files: List[Path], destination_dir: Path) -> None:
    """
    Move processed files to the designated 'processed_dir'.
    """
    for file in files:
        try:
            target_path = destination_dir / file.relative_to(to_process_dir)
            target_path.parent.mkdir(parents=True, exist_ok=True)
            file.rename(target_path)
            logging.info(f"Moved {file.name} to {target_path}")
        except Exception as e:
            logging.error(f"Failed to move file {file.name} to {target_path}: {e}")

def update_processed_files(record_path: Path, file_list: List[Path]) -> None:
    """
    Update the JSON record of processed files.
    """
    try:
        processed = set()
        if record_path.exists():
            with open(record_path, 'r') as f:
                processed = set(json.load(f))
        processed.update([file.name for file in file_list])
        with open(record_path, 'w') as f:
            json.dump(list(processed), f)
        logging.info(f"Updated processed files record at {record_path}.")
    except Exception as e:
        logging.error(f"Error updating processed files record at {record_path}: {e}")

def display_entries(entries: List[Document], file_type: str, start: int = 0, num: int = 5, text_width: int = 100, truncate: bool = False, max_length: int = 1500):
    """
    Display a subset of document entries for verification.
    """
    print(f"Displaying {num} Entries starting from index {start} in {file_type} Documents:\n")
    end = start + num
    for i, doc in enumerate(entries[start:end], start=start + 1):
        print(f"Entry {i}:")
        metadata = doc.metadata
        print(json.dumps(metadata, indent=4))
        text = doc.text
        if text:
            print("\nText:")
            if truncate and len(text) > max_length:
                displayed_text = text[:max_length] + "...\n[Text truncated]"
            else:
                displayed_text = text
            wrapped_text = textwrap.fill(displayed_text, width=text_width)
            print(wrapped_text)
        print("\n" + "-" * text_width + "\n")

########################################
# LOAD PROCESSED FILES RECORDS
########################################

# Define paths for records of processed files per file type
processed_files_record_pdf = processed_dir / "processed_files_pdf.json"
processed_files_record_excel = processed_dir / "processed_files_excel.json"
processed_files_record_ppt = processed_dir / "processed_files_ppt.json"
processed_files_record_doc = processed_dir / "processed_files_doc.json"

# Load processed files
processed_files_pdf = load_processed_files(processed_files_record_pdf)
processed_files_excel = load_processed_files(processed_files_record_excel)
processed_files_ppt = load_processed_files(processed_files_record_ppt)
processed_files_doc = load_processed_files(processed_files_record_doc)

########################################
# IDENTIFY NEW FILES FOR PROCESSING
########################################

# Supported extensions for processing
supported_extensions = [".pdf", ".xls", ".xlsx", ".ppt", ".pptx", ".doc", ".docx"]

# List all relevant files recursively from nested folders
all_files = [f for f in to_process_dir.rglob("*") if f.is_file() and f.suffix.lower() in supported_extensions]

# Categorize files by type
files_pdf = [f for f in all_files if f.suffix.lower() == ".pdf"]
files_excel = [f for f in all_files if f.suffix.lower() in [".xls", ".xlsx"]]
files_ppt = [f for f in all_files if f.suffix.lower() in [".ppt", ".pptx"]]
files_doc = [f for f in all_files if f.suffix.lower() in [".doc", ".docx"]]

# Identify new files for processing
new_files_pdf = [f for f in files_pdf if f.name not in processed_files_pdf]
new_files_excel = [f for f in files_excel if f.name not in processed_files_excel]
new_files_ppt = [f for f in files_ppt if f.name not in processed_files_ppt]
new_files_doc = [f for f in files_doc if f.name not in processed_files_doc]

# Validate new files
log_new_files(new_files_pdf, "PDF")
log_new_files(new_files_excel, "Excel")
log_new_files(new_files_ppt, "PPT")
log_new_files(new_files_doc, "DOC")

########################################
# DATA INGESTION & PREPROCESSING FUNCTIONS
########################################

def load_pdf_docs(pdf_paths: List[Path]) -> List[Document]:
    """
    Load and process PDF documents.
    """
    pdf_reader = PDFReader()
    all_docs = []
    for pdf_path in pdf_paths:
        try:
            docs = pdf_reader.load_data(pdf_path)
            docs = assign_doc_id(docs)
            for doc in docs:
                additional_info = {}
                doc = extract_metadata(doc, source_file=pdf_path.name, additional_info=additional_info)
            all_docs.extend(docs)
            logging.info(f"Processed PDF: {pdf_path.name} with {len(docs)} pages.")
        except Exception as e:
            logging.error(f"Failed to process PDF {pdf_path.name}: {e}")
    return all_docs

def load_excel_docs(excel_paths: List[Path]) -> List[Document]:
    """
    Load and process Excel documents.
    """
    all_docs = []
    for excel_path in excel_paths:
        try:
            xls = pd.ExcelFile(excel_path)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                for idx, row in df.iterrows():
                    row_dict = ", ".join(f"{col.strip().lower()}: {str(row[col]).strip()}" for col in df.columns)
                    metadata = {
                        "sheet": sheet_name,
                        "row_number": idx + 1,
                    }
                    doc = Document(text=row_dict, metadata=metadata)
                    all_docs.append(doc)
            logging.info(f"Processed Excel: {excel_path.name} with {len(df)} rows.")
        except Exception as e:
            logging.error(f"Failed to process Excel {excel_path.name}: {e}")
    all_docs = assign_doc_id(all_docs)
    for doc in all_docs:
        additional_info = {}
        doc = extract_metadata(doc, source_file=excel_path.name, additional_info=additional_info)
    return all_docs

def load_ppt_docs(ppt_paths: List[Path]) -> List[Document]:
    """
    Load and process PowerPoint documents.
    """
    all_docs = []
    for ppt_path in ppt_paths:
        try:
            prs = Presentation(ppt_path)
            for idx, slide in enumerate(prs.slides):
                slide_text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                metadata = {
                    "slide_number": idx + 1,
                }
                doc = Document(text=slide_text, metadata=metadata)
                all_docs.append(doc)
            logging.info(f"Processed PowerPoint: {ppt_path.name} with {len(prs.slides)} slides.")
        except Exception as e:
            logging.error(f"Failed to process PowerPoint {ppt_path.name}: {e}")
    all_docs = assign_doc_id(all_docs)
    for doc in all_docs:
        additional_info = {}
        doc = extract_metadata(doc, source_file=ppt_path.name, additional_info=additional_info)
    return all_docs

def split_large_section(
    section: List[str], metadata: dict, max_tokens: int, overlap_tokens: int
) -> List[Document]:
    """
    Split large text sections into manageable chunks based on token count.
    """
    all_docs = []
    text = "\n".join(section)
    tokens = tokenizer.encode(text)
    total_chunks = (len(tokens) + max_tokens - 1) // max_tokens

    for chunk_number, i in enumerate(range(0, len(tokens), max_tokens - overlap_tokens), start=1):
        chunk_tokens = tokens[i : i + max_tokens]
        chunk_text = tokenizer.decode(chunk_tokens)

        chunk_metadata = {
            **(metadata or {}),
            "chunk_number": chunk_number,
            "total_chunks_in_section": total_chunks,
            "start_token": i,
            "end_token": min(i + max_tokens, len(tokens)),
            "chunk_type": "token_chunk",
            "doc_id": str(uuid.uuid4()),
        }

        doc = Document(text=chunk_text, metadata=chunk_metadata)
        all_docs.append(doc)

    return all_docs

def chunk_by_headings(
    word_doc: WordDocument, max_tokens: int = 1024, overlap_tokens: int = 50
) -> List[Document]:
    """
    Chunk Word documents based on headings.
    """
    all_docs = []
    current_section = []
    current_metadata = None

    for paragraph in word_doc.paragraphs:
        if paragraph.style.name.startswith("Heading"):
            if current_section:
                all_docs.extend(
                    split_large_section(
                        current_section, current_metadata, max_tokens, overlap_tokens
                    )
                )
                current_section = []

            current_metadata = {"section": paragraph.text.strip()}

        current_section.append(paragraph.text.strip())

    if current_section:
        all_docs.extend(
            split_large_section(
                current_section, current_metadata, max_tokens, overlap_tokens
            )
        )

    return all_docs

def load_doc_docs(
    doc_paths: List[Path], max_tokens: int = 1024, overlap_tokens: int = 50
) -> List[Document]:
    """
    Load and process Word documents.
    """
    all_docs = []
    for doc_path in doc_paths:
        try:
            word_doc = WordDocument(doc_path)
            logging.info(f"Processing {doc_path.name} by headings.")
            docs = chunk_by_headings(word_doc, max_tokens, overlap_tokens)
            for doc in docs:
                if 'source' not in doc.metadata:
                    doc.metadata['source'] = doc_path.name
            all_docs.extend(docs)
            logging.info(f"Processed Word Document: {doc_path.name} with {len(docs)} chunks.")
        except Exception as e:
            logging.error(f"Failed to process Word Document {doc_path.name}: {e}")
    return all_docs

########################################
# PROCESSING FUNCTIONS
########################################

def trigger_processing(all_new_docs: List[Document]) -> None:
    """
    Process and index documents using Whoosh and FAISS.
    """
    if all_new_docs:
        # Create / Update Whoosh Index
        create_or_update_whoosh_index(all_new_docs, whoosh_index_path)

        # Create / Update Faiss Index
        create_or_update_faiss_index(all_new_docs, faiss_index_path, embedding_model, embedding_dimension)
    else:
        logging.info("No new documents to process.")

def scan_and_load_new_files() -> List[Document]:
    """
    Scan the 'to_process_dir' for new files and load them.
    """
    # Re-identify new files
    all_files = [f for f in to_process_dir.rglob("*") if f.is_file() and f.suffix.lower() in supported_extensions]

    # Categorize files by type
    files_pdf = [f for f in all_files if f.suffix.lower() == ".pdf"]
    files_excel = [f for f in all_files if f.suffix.lower() in [".xls", ".xlsx"]]
    files_ppt = [f for f in all_files if f.suffix.lower() in [".ppt", ".pptx"]]
    files_doc = [f for f in all_files if f.suffix.lower() in [".doc", ".docx"]]

    # Identify new files for processing
    new_files_pdf = [f for f in files_pdf if f.name not in processed_files_pdf]
    new_files_excel = [f for f in files_excel if f.name not in processed_files_excel]
    new_files_ppt = [f for f in files_ppt if f.name not in processed_files_ppt]
    new_files_doc = [f for f in files_doc if f.name not in processed_files_doc]

    # Validate new files
    log_new_files(new_files_pdf, "PDF")
    log_new_files(new_files_excel, "Excel")
    log_new_files(new_files_ppt, "PPT")
    log_new_files(new_files_doc, "DOC")

    # Load documents
    all_new_docs = []
    if new_files_pdf:
        pdf_docs = load_pdf_docs(new_files_pdf)
        logging.info(f"Total new PDF documents loaded: {len(pdf_docs)}")
        all_new_docs.extend(pdf_docs)
    if new_files_excel:
        excel_docs = load_excel_docs(new_files_excel)
        logging.info(f"Total new Excel documents loaded: {len(excel_docs)}")
        all_new_docs.extend(excel_docs)
    if new_files_ppt:
        ppt_docs = load_ppt_docs(new_files_ppt)
        logging.info(f"Total new PowerPoint documents loaded: {len(ppt_docs)}")
        all_new_docs.extend(ppt_docs)
    if new_files_doc:
        doc_docs = load_doc_docs(new_files_doc, max_tokens=max_tokens, overlap_tokens=overlap_tokens)
        logging.info(f"Total new Word documents loaded: {len(doc_docs)}")
        all_new_docs.extend(doc_docs)

    return all_new_docs

def create_or_update_whoosh_index(documents: list, index_dir: Path) -> None:
    """
    Create or update the Whoosh index with new documents.
    """
    schema = Schema(
        content=TEXT(stored=True),
        source=ID(stored=True),
        sheet=ID(stored=True),
        row_number=NUMERIC(stored=True, sortable=True),
        slide_number=NUMERIC(stored=True, sortable=True),
        section=TEXT(stored=True),
        chunk_number=NUMERIC(stored=True, sortable=True),
        total_chunks_in_section=NUMERIC(stored=True, sortable=True),
        page_number=NUMERIC(stored=True, sortable=True),
        doc_id=ID(stored=True, unique=True),
    )

    try:
        os.makedirs(index_dir, exist_ok=True)

        if not exists_in(index_dir):
            idx = create_in(index_dir, schema)
            logging.info(f"Created a new Whoosh index at {index_dir}.")
            before_count = 0
        else:
            idx = open_dir(index_dir)
            with idx.searcher() as searcher:
                before_count = searcher.doc_count()
            logging.info(f"Opened existing Whoosh index at {index_dir}. Total docs before: {before_count}.")

        writer = idx.writer()

        for doc in documents:
            try:
                content = doc.text
                source = doc.metadata.get("source", "unknown_source")
                doc_id = doc.metadata.get("doc_id", str(uuid.uuid4()))

                doc_fields = {
                    "content": content,
                    "source": source,
                    "doc_id": doc_id,
                }

                if "sheet" in doc.metadata:
                    doc_fields["sheet"] = doc.metadata["sheet"]
                if "row_number" in doc.metadata:
                    doc_fields["row_number"] = doc.metadata["row_number"]
                if "slide_number" in doc.metadata:
                    doc_fields["slide_number"] = doc.metadata["slide_number"]
                if "section" in doc.metadata:
                    doc_fields["section"] = doc.metadata["section"]
                if "chunk_number" in doc.metadata:
                    doc_fields["chunk_number"] = doc.metadata["chunk_number"]
                if "total_chunks_in_section" in doc.metadata:
                    doc_fields["total_chunks_in_section"] = doc.metadata["total_chunks_in_section"]
                if "page_number" in doc.metadata:
                    doc_fields["page_number"] = doc.metadata["page_number"]

                writer.update_document(**doc_fields)
            except Exception as e:
                logging.error(f"Failed to index document {doc.metadata.get('doc_id', 'Unknown ID')}: {e}")

        writer.commit()

        with idx.searcher() as searcher:
            after_count = searcher.doc_count()

        newly_added = len(documents)
        logging.info(
            f"Whoosh index updated. Docs before: {before_count}, docs added: {newly_added}, docs after: {after_count}."
        )
    except Exception as e:
        logging.error(f"Failed to create/update Whoosh index at {index_dir}: {e}")
        raise

def ensure_faiss_index_integrity(index_path: Path) -> bool:
    """
    Ensure that the FAISS index is readable and not corrupted.
    """
    try:
        # Attempt to read the index to check its integrity
        vector_store = FaissVectorStore.from_persist_path(str(index_path / "vector_store.faiss"))
    except RuntimeError as e:
        logging.error(f"FAISS index read error: {e}")
        return False
    return True

def create_or_update_faiss_index(
    documents: list, index_path: Path, embedding_model: NebiusEmbedding, embedding_dimension: int
) -> None:
    """
    Create or update the FAISS index with new documents.
    """
    if documents:
        newly_added = len(documents)
        if (index_path / "vector_store.faiss").exists() and ensure_faiss_index_integrity(index_path):
            # Load existing index
            vector_store = FaissVectorStore.from_persist_path(str(index_path / "vector_store.faiss"))
            storage_context = StorageContext.from_defaults(vector_store=vector_store, persist_dir=str(index_path))
            index = load_index_from_storage(storage_context, embedding=embedding_model)

            # Count vectors BEFORE
            before_count = vector_store.client.ntotal

            # Update index with new documents
            for doc in documents:
                index.insert(doc)

            # Count vectors AFTER
            after_count = vector_store.client.ntotal

            logging.info(f"Inserted {newly_added} new docs. Vectors before: {before_count}, after: {after_count}")
        else:
            # Create a new index
            vector_store = FaissVectorStore(faiss.IndexFlatIP(embedding_dimension))
            storage_context = StorageContext.from_defaults(
                vector_store=vector_store,
                docstore=SimpleDocumentStore(),
                index_store=SimpleIndexStore(),
                persist_dir=str(index_path),
            )
            index = VectorStoreIndex.from_documents(documents, storage_context=storage_context, embedding=embedding_model)

            # Count vectors BEFORE
            before_count = 0

            # Count vectors AFTER
            after_count = vector_store.client.ntotal

            logging.info(
                f"Created a new FAISS index and inserted {newly_added} docs. Vectors before: {before_count}, after: {after_count}"
            )
        # Persist the created/updated index
        index.storage_context.persist(
            persist_dir=str(index_path),
            vector_store_fname="vector_store.faiss",  # Important to save with .faiss suffix
        )
        logging.info("FAISS index created/updated and persisted successfully.")
    else:
        logging.info("No documents to index for FAISS.")

def load_all_new_documents() -> List[Document]:
    """
    Scan and load all new documents for processing.
    """
    return scan_and_load_new_files()

def move_and_record_processed_files(documents: List[Document]) -> None:
    """
    Move processed files to 'processed_dir' and update records.
    """
    # Extract file names from documents
    processed_files = set()
    for doc in documents:
        source = doc.metadata.get("source")
        if source:
            processed_files.add(source)

    # Convert to Path objects
    processed_files_paths = [to_process_dir / f for f in processed_files]

    # Move files
    move_files_to_processed(processed_files_paths, processed_dir)

    # Update processed files records based on file types
    update_processed_files(
        processed_files_record_pdf,
        [f for f in processed_files_paths if f.suffix.lower() == ".pdf"]
    )
    update_processed_files(
        processed_files_record_excel,
        [f for f in processed_files_paths if f.suffix.lower() in [".xls", ".xlsx"]]
    )
    update_processed_files(
        processed_files_record_ppt,
        [f for f in processed_files_paths if f.suffix.lower() in [".ppt", ".pptx"]]
    )
    update_processed_files(
        processed_files_record_doc,
        [f for f in processed_files_paths if f.suffix.lower() in [".doc", ".docx"]]
    )

########################################
# PROCESSING PIPELINE FUNCTION
########################################

def processing_pipeline():
    """
    The main processing pipeline to handle new documents.
    """
    # Load all new documents
    all_new_docs = load_all_new_documents()

    # Process and index documents
    trigger_processing(all_new_docs)

    # Move processed files and update records
    move_and_record_processed_files(all_new_docs)

    # Optionally, display entries for verification
    if all_new_docs:
        # Separate documents by type for display
        pdf_docs = [doc for doc in all_new_docs if doc.metadata.get("source", "").endswith(".pdf")]
        excel_docs = [doc for doc in all_new_docs if doc.metadata.get("source", "").endswith((".xls", ".xlsx"))]
        ppt_docs = [doc for doc in all_new_docs if doc.metadata.get("source", "").endswith((".ppt", ".pptx"))]
        doc_docs = [doc for doc in all_new_docs if doc.metadata.get("source", "").endswith((".doc", ".docx"))]

        # Display entries for each document type
        display_entries(pdf_docs, 'PDF')
        display_entries(excel_docs, 'Excel', start=0, num=5)
        display_entries(ppt_docs, 'PPT', start=0, num=5)
        display_entries(doc_docs, 'DOC', start=0, num=5)
    else:
        logging.info("No new documents were processed.")

########################################
# FASTAPI ENDPOINTS
########################################

@app.post("/upload")
async def upload_file(file: UploadFile):
    """
    Upload a file to the 'to_process_dir' for subsequent pipeline processing.
    """
    try:
        file_path = to_process_dir / file.filename
        with open(file_path, "wb") as f:
            f.write(await file.read())
        logging.info(f"File {file.filename} uploaded successfully to {file_path}.")
        return {"message": f"File {file.filename} uploaded to {file_path}."}
    except Exception as e:
        logging.error(f"Error uploading file: {e}")
        raise HTTPException(status_code=500, detail="File upload failed.")

@app.post("/process")
async def process_files():
    """
    Trigger document processing. Re-runs the entire pipeline on newly uploaded files.
    """
    try:
        processing_pipeline()
        return {"message": "Document processing triggered successfully."}
    except Exception as e:
        logging.error(f"Error processing documents: {e}")
        raise HTTPException(status_code=500, detail="Document processing failed.")

########################################
# GRADIO UI FUNCTIONS
########################################

def upload_file_gradio(file) -> str:
    """
    Upload file using Gradio's File object, placing it in 'to_process_dir'.
    """
    if not file:
        return "No file provided."
    try:
        file_path = to_process_dir / file.name
        with open(file_path, "wb") as f:
            f.write(file.read())
        logging.info(f"[Gradio] File {file.name} uploaded to {file_path}.")
        return f"File {file.name} uploaded to {file_path}."
    except Exception as e:
        logging.error(f"[Gradio] Error uploading file: {e}")
        return f"File upload failed: {str(e)}"

def process_files_gradio() -> str:
    """
    Trigger a re-run of the pipeline via Gradio.
    """
    try:
        processing_pipeline()
        return "Document processing triggered successfully."
    except Exception as e:
        logging.error(f"[Gradio] Error processing documents: {e}")
        return f"Document processing failed: {str(e)}"

def build_gradio_ui():
    """
    Build a Gradio Blocks UI for uploading files and re-running the pipeline.
    """
    with gr.Blocks() as demo:
        gr.Markdown("# Data Preparation Pipeline (Gradio UI)")

        with gr.Tab("Upload"):
            uploaded_file = gr.File(label="Upload a file")
            upload_output = gr.Textbox(label="Upload Status", lines=2, interactive=False)
            btn_upload = gr.Button("Upload File")
            btn_upload.click(fn=upload_file_gradio, inputs=uploaded_file, outputs=upload_output)

        with gr.Tab("Process"):
            process_output = gr.Textbox(label="Process Status", lines=2, interactive=False)
            btn_process = gr.Button("Process Documents")
            btn_process.click(fn=process_files_gradio, inputs=None, outputs=process_output)

    return demo

########################################
# GRADIO INTEGRATION WITH FASTAPI
########################################

# Define a function to launch Gradio
def run_gradio(gr_interface):
    gr_interface.launch(server_name="0.0.0.0", server_port=7860, share=False, prevent_thread_lock=True)

########################################
# RUN APPLICATION
########################################

def main():
    # Check if running in a CI/CD environment
    if os.getenv("CI") is None:
        # Launch Gradio in a separate daemon thread
        gr_interface = build_gradio_ui()
        gr_thread = threading.Thread(target=run_gradio, args=(gr_interface,), daemon=True)
        gr_thread.start()
        logging.info("Gradio interface launched.")
    else:
        logging.info("CI/CD detected. Skipping Gradio interface launch.")

    logging.info("Starting Data Preparation Pipeline.")
    # Launch FastAPI with Uvicorn
    # Only launch FastAPI if not in CI/CD, to prevent hanging
    if os.getenv("CI") is None:
        uvicorn.run(app, host="0.0.0.0", port=8080)
    else:
        # If needed, perform data preparation tasks without launching FastAPI
        processing_pipeline()
        logging.info("Data Preparation Pipeline executed successfully without launching FastAPI.")

if __name__ == "__main__":
    main()


